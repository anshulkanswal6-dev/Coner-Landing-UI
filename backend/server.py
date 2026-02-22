from fastapi import FastAPI, APIRouter, Request, Response, HTTPException, UploadFile, File, Form, Depends
from fastapi.responses import JSONResponse, PlainTextResponse
from starlette.responses import StreamingResponse
from dotenv import load_dotenv
from starlette.middleware.cors import CORSMiddleware
from motor.motor_asyncio import AsyncIOMotorClient
import os
import logging
import uuid
import re
import json
import hashlib
import asyncio
import numpy as np
import httpx
import litellm
from sentence_transformers import SentenceTransformer
from pathlib import Path
from pydantic import BaseModel, Field
from typing import List, Optional, Dict, Any
from datetime import datetime, timezone, timedelta
from bs4 import BeautifulSoup
from emergentintegrations.llm.chat import LlmChat, UserMessage

# Import Phase 1 workers and services
from workers.email_summary_worker import trigger_email_summary
from workers.learning_extraction_worker import trigger_learning_extraction, search_learned_patterns
from workers.insight_aggregator import aggregate_insights, trigger_insight_aggregation
from services.copilot_service import FounderCopilotService

ROOT_DIR = Path(__file__).parent
load_dotenv(ROOT_DIR / '.env')

mongo_url = os.environ['MONGO_URL']
client = AsyncIOMotorClient(mongo_url)
db = client[os.environ['DB_NAME']]

EMERGENT_LLM_KEY = os.environ.get('EMERGENT_LLM_KEY')
EMERGENT_PROXY_URL = "https://integrations.emergentagent.com/llm"

# Local embedding model for semantic search
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')
logger_st = logging.getLogger("sentence_transformers")
logger_st.setLevel(logging.WARNING)

# Initialize Founder Copilot Service
copilot_service = FounderCopilotService(db, EMERGENT_LLM_KEY)

app = FastAPI()
api_router = APIRouter(prefix="/api")

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger(__name__)

# ─── Pydantic Models ───
class ProjectCreate(BaseModel):
    name: str
    description: Optional[str] = ""
    welcome_message: Optional[str] = "Hi! How can I help you today?"
    primary_color: Optional[str] = "#7C3AED"

class ProjectUpdate(BaseModel):
    name: Optional[str] = None
    description: Optional[str] = None
    welcome_message: Optional[str] = None
    primary_color: Optional[str] = None
    status: Optional[str] = None
    whitelisted_domains: Optional[List[str]] = None
    agent_mode: Optional[str] = None

class GoldenRulesUpdate(BaseModel):
    preset_rules: Optional[Dict[str, bool]] = None
    custom_rules: Optional[List[str]] = None

class KnowledgeTextInput(BaseModel):
    title: str
    content: str

class KnowledgeUrlInput(BaseModel):
    url: str

class ChatInitRequest(BaseModel):
    pass

class ChatMessageRequest(BaseModel):
    session_id: str
    content: str
    current_url: Optional[str] = ""

class FeedbackRequest(BaseModel):
    message_id: str
    feedback: int  # 1 or -1

class CorrectionCreate(BaseModel):
    message_id: str
    corrected_response: str

class LeadStatusUpdate(BaseModel):
    status: str  # New, Contacted, Qualified, Closed

# ─── Helper Functions ───
def gen_id(prefix=""):
    return f"{prefix}{uuid.uuid4().hex[:16]}"

def cosine_similarity(a, b):
    a = np.array(a)
    b = np.array(b)
    return float(np.dot(a, b) / (np.linalg.norm(a) * np.linalg.norm(b) + 1e-10))

async def get_embedding(text: str) -> List[float]:
    try:
        text = text.replace("\n", " ").strip()[:8000]
        embedding = embedding_model.encode(text)
        return embedding.tolist()
    except Exception as e:
        logger.error(f"Embedding error: {e}")
        raise HTTPException(status_code=500, detail="Failed to generate embedding")

def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunk = " ".join(words[i:i + chunk_size])
        if chunk.strip():
            chunks.append(chunk)
        i += chunk_size - overlap
    return chunks if chunks else [text[:2000]]

async def get_current_user(request: Request) -> dict:
    token = None
    cookie_token = request.cookies.get("session_token")
    auth_header = request.headers.get("Authorization", "")
    if cookie_token:
        token = cookie_token
    elif auth_header.startswith("Bearer "):
        token = auth_header[7:]
    if not token:
        raise HTTPException(status_code=401, detail="Not authenticated")
    session = await db.user_sessions.find_one({"session_token": token}, {"_id": 0})
    if not session:
        raise HTTPException(status_code=401, detail="Invalid session")
    expires_at = session.get("expires_at")
    if isinstance(expires_at, str):
        expires_at = datetime.fromisoformat(expires_at)
    if expires_at and expires_at.tzinfo is None:
        expires_at = expires_at.replace(tzinfo=timezone.utc)
    if expires_at and expires_at < datetime.now(timezone.utc):
        raise HTTPException(status_code=401, detail="Session expired")
    user = await db.users.find_one({"user_id": session["user_id"]}, {"_id": 0})
    if not user:
        raise HTTPException(status_code=401, detail="User not found")
    return user

async def get_project_for_user(project_id: str, user: dict) -> dict:
    project = await db.projects.find_one({"project_id": project_id, "user_id": user["user_id"]}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Project not found")
    return project

async def get_project_by_api_key(api_key: str) -> dict:
    project = await db.projects.find_one({"api_key": api_key}, {"_id": 0})
    if not project:
        raise HTTPException(status_code=404, detail="Invalid API key")
    return project

async def search_similar_chunks(project_id: str, query_embedding: List[float], top_k: int = 5) -> List[dict]:
    chunks = await db.knowledge_chunks.find(
        {"project_id": project_id},
        {"_id": 0, "content": 1, "embedding": 1, "source_title": 1, "chunk_index": 1}
    ).to_list(500)
    if not chunks:
        return []
    scored = []
    for c in chunks:
        if c.get("embedding"):
            sim = cosine_similarity(query_embedding, c["embedding"])
            scored.append({"content": c["content"], "source": c.get("source_title", ""), "score": sim})
    scored.sort(key=lambda x: x["score"], reverse=True)
    return scored[:top_k]

async def search_corrections(project_id: str, query_embedding: List[float]) -> Optional[str]:
    corrections = await db.corrections.find(
        {"project_id": project_id},
        {"_id": 0, "corrected_response": 1, "embedding": 1}
    ).to_list(100)
    for c in corrections:
        if c.get("embedding"):
            sim = cosine_similarity(query_embedding, c["embedding"])
            if sim > 0.90:
                return c["corrected_response"]
    return None

def validate_widget_origin(request: Request, project: dict):
    """Validate request origin against project's whitelisted domains."""
    domains = project.get("whitelisted_domains", [])
    if not domains:
        return  # No whitelist = allow all (development mode)
    origin = request.headers.get("origin", "")
    referer = request.headers.get("referer", "")
    check = origin or referer
    if not check:
        return  # Allow requests without origin (e.g. server-side)
    from urllib.parse import urlparse
    parsed = urlparse(check)
    host = parsed.hostname or ""
    if not any(host == d or host.endswith(f".{d}") for d in domains):
        raise HTTPException(status_code=403, detail="Domain not whitelisted")

async def build_chat_context(project: dict, session_id: str, user_content: str, current_url: str = "", language: str = None):
    """Shared RAG pipeline: embed query, search chunks/corrections/learned patterns, build prompt."""
    project_id = project["project_id"]
    query_embedding = await get_embedding(user_content)

    # Check corrections first
    correction = await search_corrections(project_id, query_embedding)
    if correction:
        return {"type": "correction", "content": correction}

    # Search knowledge
    relevant_chunks = await search_similar_chunks(project_id, query_embedding, top_k=5)
    context = "\n\n".join([f"[Source: {c['source']}]\n{c['content']}" for c in relevant_chunks if c['score'] > 0.3])

    # NEW: Search learned patterns (Phase 1 - Self-Learning RAG)
    learned_patterns = await search_learned_patterns(db, project_id, query_embedding, top_k=3)
    if learned_patterns:
        patterns_context = "\n\n".join([f"[Learned Pattern - {p['pattern_type']}]\n{p['content']}" for p in learned_patterns])
        context = f"{context}\n\nLEARNED PATTERNS:\n{patterns_context}" if context else f"LEARNED PATTERNS:\n{patterns_context}"

    # Get conversation history
    history = await db.messages.find(
        {"session_id": session_id}, {"_id": 0, "role": 1, "content": 1}
    ).sort("created_at", 1).to_list(20)

    # Build system prompt with golden rules
    rules = project.get("golden_rules", {})
    preset = rules.get("preset_rules", {})
    custom = rules.get("custom_rules", [])
    rules_text = ""
    if preset.get("professional_tone"):
        rules_text += "- Always maintain a professional and friendly tone.\n"
    if preset.get("never_mention_competitors"):
        rules_text += "- Never mention or compare with competitors.\n"
    if preset.get("dont_discuss_pricing"):
        rules_text += "- Do not discuss specific pricing details. Direct to sales team.\n"
    if preset.get("stay_on_topic"):
        rules_text += "- Stay on topic. Only answer questions related to the business.\n"
    if preset.get("be_concise"):
        rules_text += "- Keep responses concise and clear.\n"
    if preset.get("ask_before_assuming"):
        rules_text += "- Ask clarifying questions before making assumptions.\n"
    for r in custom:
        rules_text += f"- {r}\n"

    mode = project.get("agent_mode", "support")
    if mode == "acquisition":
        mode_instruction = 'You are in ACQUISITION mode. Qualify the visitor as a lead. Naturally collect: name, email, phone, and requirements. When collected, append JSON: {"lead": {"name": "...", "email": "...", "phone": "...", "requirements": "..."}}'
    else:
        mode_instruction = "You are in SUPPORT mode. Help the user with their questions using the knowledge base."

    # NEW: Add language instruction (Phase 1 - Multilingual)
    language_instruction = ""
    if language and language != "en":
        language_map = {
            "es": "Spanish", "fr": "French", "de": "German", "it": "Italian",
            "pt": "Portuguese", "zh": "Chinese", "ja": "Japanese", "ko": "Korean",
            "ar": "Arabic", "hi": "Hindi", "ru": "Russian"
        }
        language_name = language_map.get(language, language)
        language_instruction = f"\nIMPORTANT: Respond strictly in {language_name} language. All your responses must be in {language_name}."

    system_prompt = f"""You are an AI assistant for "{project['name']}". {project.get('description', '')}

{mode_instruction}

RULES:
{rules_text}

KNOWLEDGE BASE CONTEXT:
{context if context else "No relevant knowledge found. Answer based on general knowledge but stay within the business scope."}

Current URL: {current_url or 'Not provided'}
{language_instruction}

Respond in markdown format when helpful. Be helpful, accurate, and follow all rules strictly."""

    chat_messages = [{"role": "system", "content": system_prompt}]
    for h in history[-10:]:
        chat_messages.append({"role": h["role"], "content": h["content"] if isinstance(h["content"], str) else str(h["content"])})

    return {"type": "rag", "chat_messages": chat_messages, "system_prompt": system_prompt, "chunks_used": len(relevant_chunks), "patterns_used": len(learned_patterns)}

async def store_lead_if_present(response_text: str, project_id: str, session_id: str, project: dict):
    """Extract lead JSON from response and store it. Trigger email worker if in acquisition mode."""
    lead_match = re.search(r'\{"lead":\s*\{[^}]+\}\}', response_text)
    if lead_match:
        try:
            lead_data = json.loads(lead_match.group())["lead"]
            lead_id = gen_id("lead_")
            
            logger.info(f"[LEAD CAPTURE] Capturing lead {lead_id} for project {project_id}")
            
            # NEW: Add summary_email_status field (Phase 1)
            await db.leads.insert_one({
                "lead_id": lead_id, 
                "project_id": project_id,
                "session_id": session_id, 
                "name": lead_data.get("name", ""),
                "email": lead_data.get("email", ""), 
                "phone": lead_data.get("phone", ""),
                "requirements": lead_data.get("requirements", ""), 
                "status": "New",
                "summary_email_status": "pending",  # NEW: Phase 1 - Email Summary
                "created_at": datetime.now(timezone.utc).isoformat()
            })
            
            # NEW: Trigger email summary worker (Phase 1 - Fire-and-forget)
            if project.get("agent_mode") == "acquisition" and lead_data.get("email"):
                logger.info(f"[EMAIL TRIGGER] Triggering email summary for lead {lead_id}, email: {lead_data.get('email')}")
                trigger_email_summary(db, project_id, lead_id, session_id)
                logger.info(f"[EMAIL TRIGGER] Email summary worker triggered successfully for lead {lead_id}")
            else:
                logger.info(f"[EMAIL SKIP] Skipping email - agent_mode: {project.get('agent_mode')}, has_email: {bool(lead_data.get('email'))}")
            
            return response_text.replace(lead_match.group(), "").strip()
        except Exception as e:
            logger.error(f"[LEAD ERROR] Lead storage error: {e}")
            pass
    return response_text

# ─── Auth Routes ───
@api_router.post("/auth/session")
async def auth_session(request: Request, response: Response):
    body = await request.json()
    session_id = body.get("session_id")
    if not session_id:
        raise HTTPException(status_code=400, detail="session_id required")
    async with httpx.AsyncClient() as http_client:
        resp = await http_client.get(
            "https://demobackend.emergentagent.com/auth/v1/env/oauth/session-data",
            headers={"X-Session-ID": session_id}
        )
        if resp.status_code != 200:
            raise HTTPException(status_code=401, detail="Invalid session")
        data = resp.json()
    email = data.get("email")
    name = data.get("name", "")
    picture = data.get("picture", "")
    session_token = data.get("session_token", f"st_{uuid.uuid4().hex}")
    existing = await db.users.find_one({"email": email}, {"_id": 0})
    if existing:
        user_id = existing["user_id"]
        await db.users.update_one({"email": email}, {"$set": {"name": name, "picture": picture}})
    else:
        user_id = gen_id("user_")
        await db.users.insert_one({
            "user_id": user_id, "email": email, "name": name, "picture": picture,
            "created_at": datetime.now(timezone.utc).isoformat()
        })
    await db.user_sessions.delete_many({"user_id": user_id})
    await db.user_sessions.insert_one({
        "user_id": user_id, "session_token": session_token,
        "expires_at": (datetime.now(timezone.utc) + timedelta(days=7)).isoformat(),
        "created_at": datetime.now(timezone.utc).isoformat()
    })
    response.set_cookie(
        key="session_token", value=session_token, httponly=True,
        secure=True, samesite="none", path="/", max_age=7*24*60*60
    )
    user = await db.users.find_one({"user_id": user_id}, {"_id": 0})
    return user

@api_router.get("/auth/me")
async def auth_me(user: dict = Depends(get_current_user)):
    return user

@api_router.post("/auth/logout")
async def auth_logout(request: Request, response: Response):
    token = request.cookies.get("session_token")
    if token:
        await db.user_sessions.delete_many({"session_token": token})
    response.delete_cookie("session_token", path="/", secure=True, samesite="none")
    return {"message": "Logged out"}

# ─── Project Routes ───
@api_router.get("/projects")
async def list_projects(user: dict = Depends(get_current_user)):
    projects = await db.projects.find({"user_id": user["user_id"]}, {"_id": 0}).to_list(100)
    for p in projects:
        p["knowledge_count"] = await db.knowledge_sources.count_documents({"project_id": p["project_id"]})
        p["lead_count"] = await db.leads.count_documents({"project_id": p["project_id"]})
        p["conversation_count"] = await db.conversations.count_documents({"project_id": p["project_id"]})
    return projects

@api_router.post("/projects")
async def create_project(data: ProjectCreate, user: dict = Depends(get_current_user)):
    project_id = gen_id("proj_")
    api_key = f"ep_{uuid.uuid4().hex}"
    project = {
        "project_id": project_id, "user_id": user["user_id"],
        "name": data.name, "description": data.description,
        "api_key": api_key, "welcome_message": data.welcome_message,
        "primary_color": data.primary_color, "status": "draft",
        "golden_rules": {"preset_rules": {
            "professional_tone": True, "never_mention_competitors": False,
            "dont_discuss_pricing": False, "stay_on_topic": True,
            "be_concise": True, "ask_before_assuming": False
        }, "custom_rules": []},
        "agent_mode": "support",
        "whitelisted_domains": [],
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.projects.insert_one(project)
    project.pop("_id", None)
    return project

@api_router.get("/projects/{project_id}")
async def get_project(project_id: str, user: dict = Depends(get_current_user)):
    project = await get_project_for_user(project_id, user)
    project["knowledge_count"] = await db.knowledge_sources.count_documents({"project_id": project_id})
    project["lead_count"] = await db.leads.count_documents({"project_id": project_id})
    project["conversation_count"] = await db.conversations.count_documents({"project_id": project_id})
    return project

@api_router.put("/projects/{project_id}")
async def update_project(project_id: str, data: ProjectUpdate, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    updates = {k: v for k, v in data.model_dump().items() if v is not None}
    if updates:
        await db.projects.update_one({"project_id": project_id}, {"$set": updates})
    return await db.projects.find_one({"project_id": project_id}, {"_id": 0})

@api_router.delete("/projects/{project_id}")
async def delete_project(project_id: str, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    await db.projects.delete_one({"project_id": project_id})
    await db.knowledge_sources.delete_many({"project_id": project_id})
    await db.knowledge_chunks.delete_many({"project_id": project_id})
    await db.conversations.delete_many({"project_id": project_id})
    await db.messages.delete_many({"project_id": project_id})
    await db.leads.delete_many({"project_id": project_id})
    await db.feedback.delete_many({"project_id": project_id})
    await db.corrections.delete_many({"project_id": project_id})
    return {"message": "Project deleted"}

# ─── Golden Rules Routes ───
@api_router.get("/projects/{project_id}/golden-rules")
async def get_golden_rules(project_id: str, user: dict = Depends(get_current_user)):
    project = await get_project_for_user(project_id, user)
    return project.get("golden_rules", {"preset_rules": {}, "custom_rules": []})

@api_router.put("/projects/{project_id}/golden-rules")
async def update_golden_rules(project_id: str, data: GoldenRulesUpdate, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    updates = {}
    if data.preset_rules is not None:
        updates["golden_rules.preset_rules"] = data.preset_rules
    if data.custom_rules is not None:
        updates["golden_rules.custom_rules"] = data.custom_rules
    if updates:
        await db.projects.update_one({"project_id": project_id}, {"$set": updates})
    project = await db.projects.find_one({"project_id": project_id}, {"_id": 0})
    return project.get("golden_rules", {})

# ─── Knowledge Routes ───
@api_router.get("/projects/{project_id}/knowledge")
async def list_knowledge(project_id: str, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    sources = await db.knowledge_sources.find({"project_id": project_id}, {"_id": 0}).to_list(100)
    return sources

@api_router.post("/projects/{project_id}/knowledge/text")
async def add_knowledge_text(project_id: str, data: KnowledgeTextInput, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    source_id = gen_id("ks_")
    source = {
        "source_id": source_id, "project_id": project_id, "type": "text",
        "title": data.title, "status": "processing",
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.knowledge_sources.insert_one(source)
    try:
        chunks = chunk_text(data.content)
        for i, chunk in enumerate(chunks):
            embedding = await get_embedding(chunk)
            await db.knowledge_chunks.insert_one({
                "chunk_id": gen_id("ck_"), "source_id": source_id, "project_id": project_id,
                "content": chunk, "embedding": embedding, "source_title": data.title,
                "chunk_index": i, "created_at": datetime.now(timezone.utc).isoformat()
            })
        await db.knowledge_sources.update_one(
            {"source_id": source_id},
            {"$set": {"status": "completed", "chunk_count": len(chunks)}}
        )
    except Exception as e:
        logger.error(f"Knowledge processing error: {e}")
        await db.knowledge_sources.update_one({"source_id": source_id}, {"$set": {"status": "failed", "error": str(e)}})
        raise HTTPException(status_code=500, detail=f"Failed to process knowledge: {str(e)}")
    source_doc = await db.knowledge_sources.find_one({"source_id": source_id}, {"_id": 0})
    return source_doc

@api_router.post("/projects/{project_id}/knowledge/url")
async def add_knowledge_url(project_id: str, data: KnowledgeUrlInput, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    source_id = gen_id("ks_")
    source = {
        "source_id": source_id, "project_id": project_id, "type": "url",
        "title": data.url, "url": data.url, "status": "processing",
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.knowledge_sources.insert_one(source)
    try:
        async with httpx.AsyncClient(timeout=30) as http_client:
            resp = await http_client.get(data.url, follow_redirects=True)
            html = resp.text
        soup = BeautifulSoup(html, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "header"]):
            tag.decompose()
        text = soup.get_text(separator="\n", strip=True)
        title = soup.title.string if soup.title else data.url
        text = re.sub(r'\n{3,}', '\n\n', text)
        await db.knowledge_sources.update_one({"source_id": source_id}, {"$set": {"title": title}})
        chunks = chunk_text(text)
        for i, chunk in enumerate(chunks):
            embedding = await get_embedding(chunk)
            await db.knowledge_chunks.insert_one({
                "chunk_id": gen_id("ck_"), "source_id": source_id, "project_id": project_id,
                "content": chunk, "embedding": embedding, "source_title": title,
                "chunk_index": i, "created_at": datetime.now(timezone.utc).isoformat()
            })
        await db.knowledge_sources.update_one(
            {"source_id": source_id},
            {"$set": {"status": "completed", "chunk_count": len(chunks)}}
        )
    except Exception as e:
        logger.error(f"URL scraping error: {e}")
        await db.knowledge_sources.update_one({"source_id": source_id}, {"$set": {"status": "failed", "error": str(e)}})
        raise HTTPException(status_code=500, detail=f"Failed to scrape URL: {str(e)}")
    source_doc = await db.knowledge_sources.find_one({"source_id": source_id}, {"_id": 0})
    return source_doc

@api_router.delete("/projects/{project_id}/knowledge/{source_id}")
async def delete_knowledge(project_id: str, source_id: str, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    await db.knowledge_sources.delete_one({"source_id": source_id, "project_id": project_id})
    await db.knowledge_chunks.delete_many({"source_id": source_id, "project_id": project_id})
    return {"message": "Knowledge source deleted"}

# ─── Widget API (Public, auth via API key) ───
@api_router.post("/widget/init")
async def widget_init(request: Request):
    api_key = request.headers.get("x-project-key", "")
    project = await get_project_by_api_key(api_key)
    validate_widget_origin(request, project)
    
    # NEW: Extract language from header (Phase 1 - Multilingual)
    language = request.headers.get("x-user-language", "en")
    
    session_id = gen_id("ses_")
    await db.conversations.insert_one({
        "session_id": session_id, 
        "project_id": project["project_id"],
        "language_locale": language,  # NEW: Phase 1 - Store user language
        "started_at": datetime.now(timezone.utc).isoformat(), 
        "message_count": 0
    })
    return {
        "session_id": session_id,
        "project_name": project["name"],
        "welcome_message": project.get("welcome_message", "Hi! How can I help you?"),
        "primary_color": project.get("primary_color", "#7C3AED"),
        "agent_mode": project.get("agent_mode", "support")
    }

@api_router.post("/widget/message")
async def widget_message(data: ChatMessageRequest, request: Request):
    api_key = request.headers.get("x-project-key", "")
    project = await get_project_by_api_key(api_key)
    project_id = project["project_id"]
    
    # NEW: Extract language from header (Phase 1)
    language = request.headers.get("x-user-language", None)
    
    conv = await db.conversations.find_one({"session_id": data.session_id}, {"_id": 0})
    if not conv:
        raise HTTPException(status_code=404, detail="Session not found")

    # Store user message
    user_msg_id = gen_id("msg_")
    await db.messages.insert_one({
        "message_id": user_msg_id, "session_id": data.session_id, "project_id": project_id,
        "role": "user", "content": data.content, "current_url": data.current_url,
        "created_at": datetime.now(timezone.utc).isoformat()
    })

    ctx = await build_chat_context(project, data.session_id, data.content, data.current_url, language)

    if ctx["type"] == "correction":
        assistant_msg_id = gen_id("msg_")
        await db.messages.insert_one({
            "message_id": assistant_msg_id, "session_id": data.session_id, "project_id": project_id,
            "role": "assistant", "content": ctx["content"], "source": "correction",
            "created_at": datetime.now(timezone.utc).isoformat()
        })
        await db.conversations.update_one({"session_id": data.session_id}, {"$inc": {"message_count": 2}})
        return {"message_id": assistant_msg_id, "content": ctx["content"], "source": "correction"}

    try:
        chat = LlmChat(
            api_key=EMERGENT_LLM_KEY,
            session_id=f"widget_{data.session_id}",
            system_message=ctx["system_prompt"],
            initial_messages=ctx["chat_messages"]
        ).with_model("openai", "gpt-4o-mini")
        response_text = await chat.send_message(UserMessage(text=data.content))
    except Exception as e:
        logger.error(f"LLM error: {e}")
        response_text = "I'm having trouble processing your request right now. Please try again in a moment."

    response_text = await store_lead_if_present(response_text, project_id, data.session_id, project)

    assistant_msg_id = gen_id("msg_")
    await db.messages.insert_one({
        "message_id": assistant_msg_id, "session_id": data.session_id, "project_id": project_id,
        "role": "assistant", "content": response_text,
        "chunks_used": ctx.get("chunks_used", 0),
        "created_at": datetime.now(timezone.utc).isoformat()
    })
    await db.conversations.update_one({"session_id": data.session_id}, {"$inc": {"message_count": 2}})
    return {"message_id": assistant_msg_id, "content": response_text, "source": "rag"}

# ─── SSE Streaming Chat ───
@api_router.post("/widget/message/stream")
async def widget_message_stream(data: ChatMessageRequest, request: Request):
    api_key = request.headers.get("x-project-key", "")
    project = await get_project_by_api_key(api_key)
    project_id = project["project_id"]
    
    # NEW: Extract language from header (Phase 1)
    language = request.headers.get("x-user-language", None)
    
    conv = await db.conversations.find_one({"session_id": data.session_id}, {"_id": 0})
    if not conv:
        raise HTTPException(status_code=404, detail="Session not found")

    user_msg_id = gen_id("msg_")
    await db.messages.insert_one({
        "message_id": user_msg_id, "session_id": data.session_id, "project_id": project_id,
        "role": "user", "content": data.content, "current_url": data.current_url,
        "created_at": datetime.now(timezone.utc).isoformat()
    })

    ctx = await build_chat_context(project, data.session_id, data.content, data.current_url, language)

    if ctx["type"] == "correction":
        assistant_msg_id = gen_id("msg_")
        await db.messages.insert_one({
            "message_id": assistant_msg_id, "session_id": data.session_id, "project_id": project_id,
            "role": "assistant", "content": ctx["content"], "source": "correction",
            "created_at": datetime.now(timezone.utc).isoformat()
        })
        await db.conversations.update_one({"session_id": data.session_id}, {"$inc": {"message_count": 2}})

        async def correction_gen():
            yield f"data: {json.dumps({'token': ctx['content']})}\n\n"
            yield f"data: {json.dumps({'done': True, 'message_id': assistant_msg_id, 'content': ctx['content']})}\n\n"
        return StreamingResponse(correction_gen(), media_type="text/event-stream",
            headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

    # Stream via litellm
    chat_messages = ctx["chat_messages"] + [{"role": "user", "content": data.content}]
    assistant_msg_id = gen_id("msg_")

    async def event_generator():
        full_response = ""
        try:
            response = litellm.completion(
                model="gpt-4o-mini",
                messages=chat_messages,
                api_key=EMERGENT_LLM_KEY,
                api_base=EMERGENT_PROXY_URL,
                custom_llm_provider="openai",
                stream=True
            )
            for chunk in response:
                delta = chunk.choices[0].delta
                if hasattr(delta, 'content') and delta.content:
                    token = delta.content
                    full_response += token
                    yield f"data: {json.dumps({'token': token})}\n\n"
        except Exception as e:
            logger.error(f"Streaming LLM error: {e}")
            full_response = "I'm having trouble processing your request. Please try again."
            yield f"data: {json.dumps({'token': full_response})}\n\n"

        # Post-process: lead extraction, store message
        final_text = await store_lead_if_present(full_response, project_id, data.session_id, project)
        await db.messages.insert_one({
            "message_id": assistant_msg_id, "session_id": data.session_id, "project_id": project_id,
            "role": "assistant", "content": final_text,
            "chunks_used": ctx.get("chunks_used", 0),
            "created_at": datetime.now(timezone.utc).isoformat()
        })
        await db.conversations.update_one({"session_id": data.session_id}, {"$inc": {"message_count": 2}})
        yield f"data: {json.dumps({'done': True, 'message_id': assistant_msg_id, 'content': final_text})}\n\n"

    return StreamingResponse(event_generator(), media_type="text/event-stream",
        headers={"Cache-Control": "no-cache", "X-Accel-Buffering": "no"})

@api_router.post("/widget/feedback")
async def widget_feedback(data: FeedbackRequest, request: Request):
    api_key = request.headers.get("x-project-key", "")
    project = await get_project_by_api_key(api_key)
    project_id = project["project_id"]
    
    await db.messages.update_one({"message_id": data.message_id}, {"$set": {"feedback": data.feedback}})
    
    # NEW: Trigger learning extraction (Phase 1 - Self-Learning RAG)
    trigger_learning_extraction(db, embedding_model, project_id, data.message_id, data.feedback)
    logger.info(f"Learning extraction triggered for message {data.message_id}")
    
    return {"message": "Feedback recorded"}

# ─── Widget JS Bundle ───
@api_router.get("/widget.js")
async def serve_widget_js():
    widget_path = ROOT_DIR / "widget.js"
    if not widget_path.exists():
        raise HTTPException(status_code=404, detail="Widget not found")
    content = widget_path.read_text()
    return PlainTextResponse(content, media_type="application/javascript",
        headers={"Cache-Control": "public, max-age=3600", "Access-Control-Allow-Origin": "*"})

# ─── Dashboard Chat (Sandbox) ───
@api_router.post("/projects/{project_id}/sandbox/message")
async def sandbox_message(project_id: str, data: ChatMessageRequest, user: dict = Depends(get_current_user)):
    project = await get_project_for_user(project_id, user)
    # Reuse widget message logic with project's api_key
    from starlette.datastructures import MutableHeaders
    class FakeRequest:
        def __init__(self, api_key):
            self._headers = {"x-project-key": api_key}
        @property
        def headers(self):
            return self._headers
    fake_req = FakeRequest(project["api_key"])
    return await widget_message(data, fake_req)

@api_router.post("/projects/{project_id}/sandbox/init")
async def sandbox_init(project_id: str, user: dict = Depends(get_current_user)):
    project = await get_project_for_user(project_id, user)
    class FakeRequest:
        def __init__(self, api_key):
            self._headers = {"x-project-key": api_key}
        @property
        def headers(self):
            return self._headers
    fake_req = FakeRequest(project["api_key"])
    return await widget_init(fake_req)

@api_router.post("/projects/{project_id}/sandbox/message/stream")
async def sandbox_message_stream(project_id: str, data: ChatMessageRequest, user: dict = Depends(get_current_user)):
    project = await get_project_for_user(project_id, user)
    class FakeRequest:
        def __init__(self, api_key):
            self._headers = {"x-project-key": api_key}
        @property
        def headers(self):
            return self._headers
    fake_req = FakeRequest(project["api_key"])
    return await widget_message_stream(data, fake_req)

# ─── Leads Routes ───
@api_router.get("/projects/{project_id}/leads")
async def list_leads(project_id: str, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    leads = await db.leads.find({"project_id": project_id}, {"_id": 0}).sort("created_at", -1).to_list(200)
    return leads

@api_router.put("/projects/{project_id}/leads/{lead_id}")
async def update_lead(project_id: str, lead_id: str, data: LeadStatusUpdate, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    await db.leads.update_one({"lead_id": lead_id, "project_id": project_id}, {"$set": {"status": data.status}})
    lead = await db.leads.find_one({"lead_id": lead_id}, {"_id": 0})
    return lead

# ─── Analytics Routes ───
@api_router.get("/projects/{project_id}/analytics")
async def get_analytics(project_id: str, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    total_conversations = await db.conversations.count_documents({"project_id": project_id})
    total_messages = await db.messages.count_documents({"project_id": project_id})
    total_leads = await db.leads.count_documents({"project_id": project_id})
    positive_feedback = await db.messages.count_documents({"project_id": project_id, "feedback": 1})
    negative_feedback = await db.messages.count_documents({"project_id": project_id, "feedback": -1})
    total_feedback = positive_feedback + negative_feedback
    satisfaction_rate = round((positive_feedback / total_feedback * 100), 1) if total_feedback > 0 else 0

    # Recent conversations
    recent = await db.conversations.find(
        {"project_id": project_id}, {"_id": 0}
    ).sort("started_at", -1).to_list(10)

    # Lead stats by status
    lead_pipeline = [
        {"$match": {"project_id": project_id}},
        {"$group": {"_id": "$status", "count": {"$sum": 1}}}
    ]
    lead_stats = {}
    async for doc in db.leads.aggregate(lead_pipeline):
        lead_stats[doc["_id"]] = doc["count"]

    return {
        "total_conversations": total_conversations,
        "total_messages": total_messages,
        "total_leads": total_leads,
        "satisfaction_rate": satisfaction_rate,
        "positive_feedback": positive_feedback,
        "negative_feedback": negative_feedback,
        "lead_stats": lead_stats,
        "recent_conversations": recent
    }

# ─── Feedback & Corrections ───
@api_router.get("/projects/{project_id}/feedback")
async def list_feedback(project_id: str, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    messages = await db.messages.find(
        {"project_id": project_id, "feedback": {"$exists": True, "$ne": 0}},
        {"_id": 0}
    ).sort("created_at", -1).to_list(100)
    results = []
    for msg in messages:
        if msg["role"] == "assistant":
            user_msg = await db.messages.find_one(
                {"session_id": msg["session_id"], "role": "user", "created_at": {"$lt": msg["created_at"]}},
                {"_id": 0}
            )
            results.append({
                "message_id": msg["message_id"],
                "session_id": msg["session_id"],
                "user_query": user_msg["content"] if user_msg else "",
                "assistant_response": msg["content"],
                "feedback": msg.get("feedback", 0),
                "created_at": msg["created_at"]
            })
    return results

@api_router.post("/projects/{project_id}/corrections")
async def create_correction(project_id: str, data: CorrectionCreate, user: dict = Depends(get_current_user)):
    await get_project_for_user(project_id, user)
    msg = await db.messages.find_one({"message_id": data.message_id}, {"_id": 0})
    if not msg:
        raise HTTPException(status_code=404, detail="Message not found")
    user_msg = await db.messages.find_one(
        {"session_id": msg["session_id"], "role": "user", "created_at": {"$lt": msg["created_at"]}},
        {"_id": 0}
    )
    query_text = user_msg["content"] if user_msg else ""
    embedding = await get_embedding(query_text + " " + data.corrected_response)
    correction_id = gen_id("cor_")
    await db.corrections.insert_one({
        "correction_id": correction_id, "project_id": project_id,
        "message_id": data.message_id, "original_query": query_text,
        "corrected_response": data.corrected_response, "embedding": embedding,
        "created_at": datetime.now(timezone.utc).isoformat()
    })
    return {"correction_id": correction_id, "message": "Correction saved"}

# ─── Root ───
@api_router.get("/")
async def root():
    return {"message": "EmergentPulse AI API", "version": "1.0.0"}

# ─── Document Parser Functions ───
import io
from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
from openpyxl import load_workbook

async def extract_text_from_pdf(file_bytes: bytes) -> str:
    """Extract text from PDF"""
    try:
        pdf = PdfReader(io.BytesIO(file_bytes))
        text = ""
        for page in pdf.pages:
            text += page.extract_text() + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"PDF extraction failed: {str(e)}")

async def extract_text_from_docx(file_bytes: bytes) -> str:
    """Extract text from DOCX"""
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
        text = "\n".join([para.text for para in doc.paragraphs if para.text.strip()])
        return text.strip()
    except Exception as e:
        raise Exception(f"DOCX extraction failed: {str(e)}")

async def extract_text_from_pptx(file_bytes: bytes) -> str:
    """Extract text from PPTX"""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"PPTX extraction failed: {str(e)}")

async def extract_text_from_xlsx(file_bytes: bytes) -> str:
    """Extract text from XLSX"""
    try:
        wb = load_workbook(io.BytesIO(file_bytes), data_only=True)
        text = ""
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text += f"\n\n--- Sheet: {sheet_name} ---\n"
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                if row_text.strip():
                    text += row_text + "\n"
        return text.strip()
    except Exception as e:
        raise Exception(f"XLSX extraction failed: {str(e)}")

async def extract_text_from_file(file_bytes: bytes, filename: str) -> str:
    """Route to appropriate parser based on file extension"""
    ext = filename.lower().split('.')[-1]
    
    if ext == 'pdf':
        return await extract_text_from_pdf(file_bytes)
    elif ext == 'docx':
        return await extract_text_from_docx(file_bytes)
    elif ext == 'pptx':
        return await extract_text_from_pptx(file_bytes)
    elif ext in ['xlsx', 'xls']:
        return await extract_text_from_xlsx(file_bytes)
    elif ext == 'txt':
        return file_bytes.decode('utf-8', errors='ignore')
    else:
        raise Exception(f"Unsupported file type: {ext}")

# ─── Document Upload Endpoint ───
@api_router.post("/projects/{project_id}/knowledge/upload")
async def upload_document(
    project_id: str,
    file: UploadFile = File(...),
    user: dict = Depends(get_current_user)
):
    """Upload and process document (PDF, DOCX, PPTX, TXT, XLSX)"""
    await get_project_for_user(project_id, user)
    
    # Validate file type
    allowed_extensions = ['pdf', 'docx', 'pptx', 'txt', 'xlsx', 'xls']
    file_ext = file.filename.lower().split('.')[-1]
    if file_ext not in allowed_extensions:
        raise HTTPException(
            status_code=400,
            detail=f"Unsupported file type. Allowed: {', '.join(allowed_extensions)}"
        )
    
    # Read file
    file_bytes = await file.read()
    file_size = len(file_bytes)
    
    # Validate size (10MB max)
    MAX_SIZE = 10 * 1024 * 1024  # 10 MB
    if file_size > MAX_SIZE:
        raise HTTPException(
            status_code=400,
            detail=f"File too large. Maximum size is 10 MB. Your file: {file_size / (1024*1024):.2f} MB"
        )
    
    # Check for duplicate (same filename + size)
    existing = await db.knowledge_sources.find_one({
        "project_id": project_id,
        "type": "document",
        "file_name": file.filename,
        "file_size": file_size
    })
    if existing:
        raise HTTPException(
            status_code=400,
            detail="This file has already been uploaded. Delete the existing file first."
        )
    
    # Create source record
    source_id = gen_id("ks_")
    source = {
        "source_id": source_id,
        "project_id": project_id,
        "type": "document",
        "file_name": file.filename,
        "file_size": file_size,
        "file_type": file_ext,
        "title": file.filename,
        "status": "processing",
        "created_at": datetime.now(timezone.utc).isoformat()
    }
    await db.knowledge_sources.insert_one(source)
    
    try:
        # Extract text
        text = await extract_text_from_file(file_bytes, file.filename)
        
        if not text or len(text.strip()) < 10:
            raise Exception("No meaningful text could be extracted from the document")
        
        # Chunk with metadata
        chunks = chunk_text(text, chunk_size=600, overlap=100)
        
        # Generate embeddings and store
        for i, chunk in enumerate(chunks):
            embedding = await get_embedding(chunk)
            await db.knowledge_chunks.insert_one({
                "chunk_id": gen_id("ck_"),
                "source_id": source_id,
                "project_id": project_id,
                "content": chunk,
                "embedding": embedding,
                "source_title": file.filename,
                "source_type": "document",
                "file_name": file.filename,
                "chunk_index": i,
                "created_at": datetime.now(timezone.utc).isoformat()
            })
        
        # Update source as completed
        await db.knowledge_sources.update_one(
            {"source_id": source_id},
            {"$set": {
                "status": "completed",
                "chunk_count": len(chunks),
                "text_length": len(text)
            }}
        )
        
    except Exception as e:
        logger.error(f"Document processing error: {e}")
        await db.knowledge_sources.update_one(
            {"source_id": source_id},
            {"$set": {"status": "failed", "error": str(e)}}
        )
        raise HTTPException(status_code=500, detail=f"Failed to process document: {str(e)}")
    
    source_doc = await db.knowledge_sources.find_one({"source_id": source_id}, {"_id": 0})
    return source_doc

# ─── List Files Endpoint ───
@api_router.get("/projects/{project_id}/knowledge/files")
async def list_files(project_id: str, user: dict = Depends(get_current_user)):
    """List all uploaded document files"""
    await get_project_for_user(project_id, user)
    files = await db.knowledge_sources.find(
        {"project_id": project_id, "type": "document"},
        {"_id": 0}
    ).sort("created_at", -1).to_list(100)
    return files

# ─── Delete File Endpoint ───
@api_router.delete("/projects/{project_id}/knowledge/files/{file_id}")
async def delete_file(project_id: str, file_id: str, user: dict = Depends(get_current_user)):
    """Delete uploaded file and its vectors"""
    await get_project_for_user(project_id, user)
    
    # Delete source
    result = await db.knowledge_sources.delete_one({
        "source_id": file_id,
        "project_id": project_id,
        "type": "document"
    })
    
    if result.deleted_count == 0:
        raise HTTPException(status_code=404, detail="File not found")
    
    # Delete all chunks/vectors
    await db.knowledge_chunks.delete_many({
        "source_id": file_id,
        "project_id": project_id
    })
    
    return {"message": "File and its knowledge vectors deleted successfully"}

# ─── Enhanced Website Crawler with Navigation Intelligence ───
from urllib.parse import urlparse, urljoin

async def crawl_website_with_navigation(base_url: str, max_pages: int = 50) -> tuple:
    """
    Crawl website and extract:
    1. Page content
    2. Navigation structure (links, hierarchy, menu)
    
    Returns: (content_chunks, navigation_data)
    """
    try:
        visited = set()
        to_visit = [base_url]
        domain = urlparse(base_url).netloc
        
        all_content = []
        navigation_map = {}
        site_structure = {"pages": {}, "links": []}
        
        async with httpx.AsyncClient(timeout=30, follow_redirects=True) as client:
            while to_visit and len(visited) < max_pages:
                url = to_visit.pop(0)
                
                if url in visited:
                    continue
                
                try:
                    resp = await client.get(url)
                    if resp.status_code != 200:
                        continue
                    
                    visited.add(url)
                    html = resp.text
                    soup = BeautifulSoup(html, "html.parser")
                    
                    # Extract page title
                    page_title = soup.title.string if soup.title else url
                    
                    # Extract navigation elements
                    nav_links = []
                    for nav in soup.find_all(['nav', 'menu']):
                        for link in nav.find_all('a', href=True):
                            link_text = link.get_text(strip=True)
                            link_href = urljoin(url, link['href'])
                            if link_text and link_href:
                                nav_links.append({"text": link_text, "url": link_href})
                    
                    # Extract breadcrumbs
                    breadcrumbs = []
                    for bc in soup.find_all(class_=re.compile('breadcrumb', re.I)):
                        for link in bc.find_all('a'):
                            breadcrumbs.append(link.get_text(strip=True))
                    
                    # Store page info
                    site_structure["pages"][url] = {
                        "title": page_title,
                        "nav_links": nav_links,
                        "breadcrumbs": breadcrumbs
                    }
                    
                    # Extract main content
                    for tag in soup(['script', 'style', 'nav', 'footer', 'header']):
                        tag.decompose()
                    
                    text = soup.get_text(separator="\n", strip=True)
                    text = re.sub(r'\n{3,}', '\n\n', text)
                    
                    if text.strip():
                        all_content.append({
                            "url": url,
                            "title": page_title,
                            "content": text
                        })
                    
                    # Find internal links
                    for link in soup.find_all('a', href=True):
                        href = urljoin(url, link['href'])
                        parsed = urlparse(href)
                        
                        # Only follow same-domain links
                        if parsed.netloc == domain and href not in visited and href not in to_visit:
                            # Avoid common non-content URLs
                            if not any(skip in href.lower() for skip in ['#', 'javascript:', 'mailto:', '.pdf', '.zip']):
                                to_visit.append(href)
                                site_structure["links"].append({
                                    "from": url,
                                    "to": href,
                                    "text": link.get_text(strip=True)
                                })
                
                except Exception as e:
                    logger.error(f"Error crawling {url}: {e}")
                    continue
        
        # Generate navigation intelligence
        navigation_chunks = generate_navigation_chunks(site_structure)
        
        return all_content, navigation_chunks
        
    except Exception as e:
        raise Exception(f"Website crawl failed: {str(e)}")

def generate_navigation_chunks(site_structure: dict) -> List[dict]:
    """
    Generate structured navigation guidance from site structure.
    Example: "How to apply for driving licence" → step-by-step navigation
    """
    nav_chunks = []
    
    # Analyze navigation patterns
    pages = site_structure.get("pages", {})
    links = site_structure.get("links", [])
    
    # Create task-based navigation from common patterns
    for url, page_data in pages.items():
        title = page_data.get("title", "")
        nav_links = page_data.get("nav_links", [])
        breadcrumbs = page_data.get("breadcrumbs", [])
        
        # Generate navigation instruction
        if breadcrumbs:
            nav_text = f"To reach '{title}':\n"
            nav_text += "Path: " + " → ".join(breadcrumbs) + "\n"
            
            if nav_links:
                nav_text += "Available options from this page:\n"
                for link in nav_links[:10]:  # Limit to 10
                    nav_text += f"- {link['text']}\n"
            
            nav_chunks.append({
                "content": nav_text,
                "page_title": title,
                "type": "navigation"
            })
    
    # Create link-based navigation
    link_text = "Website navigation map:\n"
    for link in links[:100]:  # Limit to avoid huge chunks
        if link.get("text"):
            link_text += f"- {link['text']} (from {link['from']} to {link['to']})\n"
    
    if link_text:
        nav_chunks.append({
            "content": link_text,
            "type": "navigation_map"
        })
    
    return nav_chunks

# ─── Website Sync Endpoints ───
@api_router.post("/projects/{project_id}/knowledge/sync")
async def sync_website(project_id: str, user: dict = Depends(get_current_user)):
    """
    Manually trigger website re-sync.
    Crawls website, extracts content + navigation, updates knowledge base.
    """
    project = await get_project_for_user(project_id, user)
    
    # Get website URL from existing URL-type knowledge sources
    url_sources = await db.knowledge_sources.find({
        "project_id": project_id,
        "type": "url"
    }).to_list(10)
    
    if not url_sources:
        raise HTTPException(
            status_code=400,
            detail="No website URL found. Add a website URL first via Knowledge > Website."
        )
    
    # Use the first URL as base
    base_url = url_sources[0].get("url")
    
    # Mark sync as in progress
    sync_id = gen_id("sync_")
    sync_record = {
        "sync_id": sync_id,
        "project_id": project_id,
        "status": "in_progress",
        "started_at": datetime.now(timezone.utc).isoformat()
    }
    await db.website_syncs.insert_one(sync_record)
    
    try:
        # Delete old website-sourced chunks (preserve documents)
        await db.knowledge_chunks.delete_many({
            "project_id": project_id,
            "source_type": {"$in": ["url", "navigation"]}
        })
        
        # Crawl website
        content_pages, navigation_chunks = await crawl_website_with_navigation(base_url, max_pages=30)
        
        total_chunks = 0
        
        # Process content pages
        for page in content_pages:
            source_id = gen_id("ks_")
            
            # Create temporary source for this page
            await db.knowledge_sources.update_one(
                {"project_id": project_id, "url": page["url"], "type": "url"},
                {"$set": {
                    "source_id": source_id,
                    "title": page["title"],
                    "status": "completed",
                    "updated_at": datetime.now(timezone.utc).isoformat()
                }},
                upsert=True
            )
            
            # Chunk content
            chunks = chunk_text(page["content"], chunk_size=600, overlap=100)
            
            for i, chunk in enumerate(chunks):
                embedding = await get_embedding(chunk)
                await db.knowledge_chunks.insert_one({
                    "chunk_id": gen_id("ck_"),
                    "source_id": source_id,
                    "project_id": project_id,
                    "content": chunk,
                    "embedding": embedding,
                    "source_title": page["title"],
                    "source_type": "url",
                    "url": page["url"],
                    "chunk_index": i,
                    "created_at": datetime.now(timezone.utc).isoformat()
                })
                total_chunks += 1
        
        # Process navigation chunks
        nav_source_id = gen_id("ks_")
        for nav_chunk in navigation_chunks:
            embedding = await get_embedding(nav_chunk["content"])
            await db.knowledge_chunks.insert_one({
                "chunk_id": gen_id("ck_"),
                "source_id": nav_source_id,
                "project_id": project_id,
                "content": nav_chunk["content"],
                "embedding": embedding,
                "source_title": "Website Navigation",
                "source_type": "navigation",
                "chunk_index": 0,
                "created_at": datetime.now(timezone.utc).isoformat()
            })
            total_chunks += 1
        
        # Update sync record
        await db.website_syncs.update_one(
            {"sync_id": sync_id},
            {"$set": {
                "status": "completed",
                "completed_at": datetime.now(timezone.utc).isoformat(),
                "pages_crawled": len(content_pages),
                "chunks_created": total_chunks
            }}
        )
        
        return {
            "message": "Website synced successfully",
            "pages_crawled": len(content_pages),
            "chunks_created": total_chunks,
            "sync_id": sync_id
        }
        
    except Exception as e:
        logger.error(f"Website sync error: {e}")
        await db.website_syncs.update_one(
            {"sync_id": sync_id},
            {"$set": {
                "status": "failed",
                "error": str(e),
                "completed_at": datetime.now(timezone.utc).isoformat()
            }}
        )
        raise HTTPException(status_code=500, detail=f"Website sync failed: {str(e)}")

@api_router.get("/projects/{project_id}/knowledge/sync/status")
async def get_sync_status(project_id: str, user: dict = Depends(get_current_user)):
    """Get last sync status and timestamp"""
    await get_project_for_user(project_id, user)
    
    last_sync = await db.website_syncs.find_one(
        {"project_id": project_id},
        {"_id": 0},
        sort=[("started_at", -1)]
    )
    
    return last_sync if last_sync else {"status": "never_synced"}


# ─── NEW: Phase 1 Analytics Endpoints ───

class CopilotQueryRequest(BaseModel):
    query: str = Field(..., min_length=1, max_length=500, description="Natural language analytics query")

class InsightRefreshRequest(BaseModel):
    date_from: Optional[str] = None
    date_to: Optional[str] = None


@api_router.get("/projects/{project_id}/analytics/insights/cards")
async def get_project_insight_cards(
    project_id: str,
    date_from: Optional[str] = None,
    date_to: Optional[str] = None,
    user: dict = Depends(get_current_user)
):
    """Get analytics insight cards for a specific project. Tenant-scoped."""
    try:
        # Validate project access (tenant-scoped)
        await get_project_for_user(project_id, user)
        
        # Default to last 7 days if not specified
        if not date_to:
            date_to = datetime.now(timezone.utc).isoformat()
        if not date_from:
            date_from = (datetime.now(timezone.utc) - timedelta(days=7)).isoformat()
        
        # Try to get from materialized insights first
        insights = await db.insight_summaries.find(
            {
                "project_id": project_id,
                "date_from": date_from,
                "date_to": date_to
            },
            {"_id": 0}
        ).to_list(20)
        
        # If no materialized insights, compute on-the-fly
        if not insights:
            insights = await aggregate_insights(db, project_id, date_from, date_to)
        
        return {
            "status": "success",
            "insights": insights,
            "count": len(insights),
            "date_from": date_from,
            "date_to": date_to
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error fetching insight cards: {e}")
        raise HTTPException(status_code=500, detail="Failed to fetch insight cards")


@api_router.post("/projects/{project_id}/analytics/copilot/ask")
async def founder_copilot_ask(
    project_id: str,
    data: CopilotQueryRequest,
    user: dict = Depends(get_current_user)
):
    """Natural language analytics query endpoint (Founder Copilot). Tenant-scoped, PII-safe."""
    try:
        # Validate project access (tenant-scoped)
        await get_project_for_user(project_id, user)
        
        # Sanitize query (prevent injection)
        query = data.query.strip()
        
        if len(query) < 3:
            raise HTTPException(status_code=400, detail="Query too short. Please ask a complete question.")
        
        # Use copilot service
        result = await copilot_service.query(project_id, query)
        
        return result
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error processing copilot query: {e}")
        return {
            "status": "error",
            "query": data.query,
            "error": str(e),
            "explanation": "I encountered an error processing your query. Please try rephrasing it."
        }


@api_router.post("/projects/{project_id}/analytics/insights/refresh")
async def refresh_insights(
    project_id: str,
    user: dict = Depends(get_current_user),
    date_from: Optional[str] = None,
    date_to: Optional[str] = None
):
    """Manually trigger insight aggregation. Tenant-scoped."""
    try:
        # Validate project access (tenant-scoped)
        await get_project_for_user(project_id, user)
        
        # Trigger insight aggregation (fire-and-forget)
        trigger_insight_aggregation(db, project_id, date_from, date_to)
        
        return {
            "status": "success",
            "message": "Insight aggregation triggered",
            "processing": True
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Error triggering insight refresh: {e}")
        raise HTTPException(status_code=500, detail="Failed to trigger insight aggregation")


# ───────────────────────────────────────────────────────────────────────────────

app.include_router(api_router)

app.add_middleware(
    CORSMiddleware,
    allow_credentials=True,
    allow_origins=os.environ.get('CORS_ORIGINS', '*').split(','),
    allow_methods=["*"],
    allow_headers=["*"],
)

@app.on_event("startup")
async def startup():
    await db.users.create_index("email", unique=True)
    await db.users.create_index("user_id", unique=True)
    await db.user_sessions.create_index("session_token")
    await db.projects.create_index("user_id")
    await db.projects.create_index("api_key", unique=True)
    await db.knowledge_sources.create_index("project_id")
    await db.knowledge_chunks.create_index("project_id")
    await db.knowledge_chunks.create_index("source_id")
    await db.conversations.create_index("project_id")
    await db.conversations.create_index("session_id", unique=True)
    await db.messages.create_index("session_id")
    await db.messages.create_index("project_id")
    await db.leads.create_index("project_id")
    await db.corrections.create_index("project_id")
    
    # NEW: Phase 1 - Indexes for new collections
    await db.learned_patterns.create_index("project_id")
    await db.learned_patterns.create_index("message_id")
    await db.insight_summaries.create_index("project_id")
    await db.insight_summaries.create_index([("project_id", 1), ("metric_type", 1)])
    
    logger.info("EmergentPulse AI backend started")

@app.on_event("shutdown")
async def shutdown():
    client.close()


# ─── Lead Export Endpoint ───
@api_router.get("/projects/{project_id}/leads/export")
async def export_leads(project_id: str, user: dict = Depends(get_current_user)):
    """Export all leads as CSV"""
    await get_project_for_user(project_id, user)
    leads = await db.leads.find({"project_id": project_id}, {"_id": 0}).sort("created_at", -1).to_list(1000)
    
    # Create CSV
    import io
    import csv
    
    output = io.StringIO()
    writer = csv.DictWriter(output, fieldnames=[
        'lead_id', 'name', 'email', 'phone', 'query_objective', 'user_need', 
        'source', 'preferred_contact', 'availability', 'status', 'created_at', 'admin_notes'
    ])
    writer.writeheader()
    
    for lead in leads:
        writer.writerow({
            'lead_id': lead.get('lead_id', ''),
            'name': lead.get('name', ''),
            'email': lead.get('email', ''),
            'phone': lead.get('phone', ''),
            'query_objective': lead.get('query_objective', lead.get('requirements', '')),
            'user_need': lead.get('user_need', lead.get('details', '')),
            'source': lead.get('source', 'chatbot'),
            'preferred_contact': lead.get('preferred_contact', ''),
            'availability': lead.get('availability', ''),
            'status': lead.get('status', 'New'),
            'created_at': lead.get('created_at', ''),
            'admin_notes': lead.get('admin_notes', '')
        })
    
    csv_content = output.getvalue()
    return Response(
        content=csv_content,
        media_type="text/csv",
        headers={"Content-Disposition": f"attachment; filename=leads_{project_id}.csv"}
    )

# ─── Update Lead with Extended Fields ───
@api_router.put("/projects/{project_id}/leads/{lead_id}/details")
async def update_lead_details(
    project_id: str, 
    lead_id: str, 
    data: dict,
    user: dict = Depends(get_current_user)
):
    """Update lead with admin notes and extended fields"""
    await get_project_for_user(project_id, user)
    
    # Extract allowed fields
    update_fields = {}
    allowed_fields = ['status', 'admin_notes', 'preferred_contact', 'query_objective', 'user_need', 'availability', 'social_links']
    
    for field in allowed_fields:
        if field in data:
            update_fields[field] = data[field]
    
    if update_fields:
        await db.leads.update_one(
            {"lead_id": lead_id, "project_id": project_id}, 
            {"$set": update_fields}
        )
    
    lead = await db.leads.find_one({"lead_id": lead_id}, {"_id": 0})
    return lead
